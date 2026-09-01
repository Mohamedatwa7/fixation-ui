"""Next Steps pillar on an empty slide — matches the black-bar + bordered-box
pillar format of the Training/Validation slide so it can be copy-pasted in as
a third column. -> Downloads/F1X8_NextSteps_Pillar.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1A, 0x1A, 0x2E)
LINE = RGBColor(0xBF, 0xBF, 0xBF)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2F, 0x4B, 0xE0)
GREEN = RGBColor(0x1E, 0x9E, 0x4A)
ORANGE = RGBColor(0xE8, 0x86, 0x2B)

# (number+lead [underlined], [(text, color, bold), ...] rest, timing)
ITEMS = [
    ("1. WEEKLY LEARNING",
     [(" - Model retrains on every new campaign's results; ", NAVY, False),
      ("accuracy compounds automatically", BLUE, True),
      (" with zero manual work", NAVY, False)],
     "live this week"),
    ("2. VIDEO UPGRADE",
     [(" - Bring reels scoring to image-level precision: ", NAVY, False),
      ("73% → 85% accuracy", GREEN, True),
      (" for the majority of our content", NAVY, False)],
     "2–3 weeks"),
    ("3. PAID CALIBRATION",
     [(" - Scores tuned to ", NAVY, False),
      ("CTR & cost per engagement", BLUE, True),
      (" using ads-manager reports (FF8 showed a 3.5× spread between sister assets)", NAVY, False)],
     "next campaigns"),
    ("4. SCALE",
     [(" - Onboard more Samsung accounts (", NAVY, False),
      ("12.5K posts ready", BLUE, True),
      (") and make the pre-flight check standard before every boost", NAVY, False)],
     "1–2 months"),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

# Pillar sized to sit as a third column beside the two existing pillars.
PX, PY = Inches(4.67), Inches(1.55)
PW, PH = Inches(4.0), Inches(4.35)
BAR_H = Inches(0.38)

bar = slide.shapes.add_shape(1, PX, PY, PW, BAR_H)
bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
bar.line.fill.background()
tf = bar.text_frame
tf.margin_top = Emu(18288); tf.margin_bottom = Emu(18288)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Next Steps & Scaling"
r.font.size = Pt(14); r.font.bold = True; r.font.name = "Segoe UI"; r.font.color.rgb = WHITE

box = slide.shapes.add_shape(1, PX, Emu(int(PY) + int(BAR_H)), PW, Emu(int(PH) - int(BAR_H)))
box.fill.solid(); box.fill.fore_color.rgb = WHITE
box.line.color.rgb = LINE; box.line.width = Emu(9525)
box.shadow.inherit = False
tf = box.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
tf.margin_left = Emu(118872); tf.margin_right = Emu(118872); tf.margin_top = Emu(100584)

first = True
for lead, rest_parts, timing in ITEMS:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(9)
    r = p.add_run()
    r.text = lead
    r.font.size = Pt(11); r.font.bold = True; r.font.underline = True
    r.font.name = "Segoe UI"; r.font.color.rgb = NAVY
    for text, color, bold in rest_parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(11); r.font.bold = bold
        r.font.name = "Segoe UI"; r.font.color.rgb = color
    r = p.add_run()
    r.text = f"  ({timing})"
    r.font.size = Pt(10.5); r.font.bold = True
    r.font.name = "Segoe UI"; r.font.color.rgb = ORANGE

p = tf.add_paragraph()
p.alignment = PP_ALIGN.LEFT
p.space_before = Pt(4)
r = p.add_run()
r.text = "Every campaign scored makes the model smarter — the earlier it spreads, the faster it compounds."
r.font.size = Pt(10.5); r.font.bold = True; r.font.italic = True
r.font.name = "Segoe UI"; r.font.color.rgb = BLUE

out = os.path.join(os.environ.get("USERPROFILE", ""), "Downloads", "F1X8_NextSteps_Pillar.pptx")
prs.save(out)
print("wrote", out)
