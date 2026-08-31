"""Generate the TRIBE-tool vs F1X8 comparison slide (16:9, noir brand theme).
Usage: python docs/make_tribe_slide.py  ->  docs/TRIBE_vs_F1X8.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

NOIR = RGBColor(0x0A, 0x0A, 0x0A)
PANEL = RGBColor(0x14, 0x14, 0x14)
ELEV = RGBColor(0x1C, 0x1C, 0x1C)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)
MUTED = RGBColor(0x9A, 0x9A, 0x9A)
ACCENT = RGBColor(0xE0, 0xE0, 0xE0)
GREEN = RGBColor(0x8A, 0xE1, 0xB8)

ROWS = [
    ("The question it answers",
     "“How would a human brain react to this content?”",
     "“How will this ad perform for our brand — and how do we fix it?”"),
    ("What it's built on",
     "Meta's research model — brain scans of volunteers watching movies & podcasts",
     "Real eye-tracking data + Samsung Gulf's own posts and their real engagement results"),
    ("Tuned to our brand?",
     "General-purpose — same prediction for any company's content",
     "Calibrated on our own feed; scores every creative against our own past posts"),
    ("Follows brand guidelines?",
     "Not brand-aware",
     "Samsung Brand Creative Playbook built into every score and recommendation"),
    ("What you get back",
     "A prediction of neural response",
     "Performance score + specific risks + shoot-ready revision brief + AI-revised creative"),
    ("Proven against real results?",
     "Validated on brain activity — not yet on ad performance",
     "Picks the better performer ~85 times out of 100 on held-out Samsung posts (coin flip = 50)"),
    ("Combining them — tested?",
     "—",
     "Tested on 171 real posts: adding brain predictions did not improve accuracy — F1X8 already captures it"),
]

FOOTER = ("“It predicts what a brain does; F1X8 predicts what our feed does — "
          "verified against our own real results.”")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = NOIR


def textbox(x, y, w, h, text, size, color, bold=False, font="Segoe UI",
            align=PP_ALIGN.LEFT, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = color
    return tb


# Header
textbox(Inches(0.6), Inches(0.32), Inches(9), Inches(0.5),
        "BRAIN-RESPONSE TOOL (TRIBE)  VS  F1X8", 13, MUTED, font="Consolas")
textbox(Inches(0.6), Inches(0.72), Inches(11), Inches(0.7),
        "Two tools, two different questions", 30, WHITE, bold=True)

# Table
left, top = Inches(0.6), Inches(1.62)
width, height = Inches(12.13), Inches(4.9)
n_rows = len(ROWS) + 1
tbl_shape = slide.shapes.add_table(n_rows, 3, left, top, width, height)
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.55)
tbl.columns[1].width = Inches(4.29)
tbl.columns[2].width = Inches(5.29)

headers = ("", "Brain-response tool (built on Meta TRIBE)", "F1X8")
for c, htext in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ELEV
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = htext
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Consolas"
    r.font.color.rgb = ACCENT if c != 2 else GREEN

for i, (label, tribe, f1x8) in enumerate(ROWS, start=1):
    for c, (text, color, bold, size) in enumerate((
        (label, MUTED, True, 11),
        (tribe, WHITE, False, 11.5),
        (f1x8, WHITE, False, 11.5),
    )):
        cell = tbl.cell(i, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if i % 2 else NOIR
        cell.margin_left = Emu(91440)
        cell.margin_right = Emu(91440)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.name = "Consolas" if c == 0 else "Segoe UI"
        f.color.rgb = GREEN if (c == 2 and i >= len(ROWS) - 1) else color

# Footer one-liner
textbox(Inches(0.6), Inches(6.72), Inches(12.13), Inches(0.6),
        FOOTER, 14, ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Speaker notes
notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Positioning: do not frame as competition. The TRIBE-based tool is genuine "
    "cutting-edge research; the two answer different questions. TRIBE predicts how an "
    "average brain responds to content; F1X8 predicts and diagnoses feed performance "
    "for our brand specifically.\n\n"
    "The combination question ('why not use both?') has an evidence-based answer: we ran "
    "Meta's model in its full licensed configuration over 171 real Samsung posts with "
    "known engagement outcomes. Its brain-response features added zero predictive "
    "accuracy on top of F1X8 (best case +0.008 AUC on images; on videos it reduced "
    "accuracy), while adding GPU cost and minutes of processing per creative. Reason: "
    "feed performance is driven by behavior (stopping the scroll, watching, sharing), "
    "which F1X8 measures directly against real results — the brain prediction has "
    "nothing left to add.\n\n"
    "Accuracy claim source: fine-tuned ranker, holdout AUC 0.851 on unseen Samsung "
    "creatives (~85/100 pairwise). Combination test: eval/tribe/README.md in the repo."
)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TRIBE_vs_F1X8.pptx")
prs.save(out)
print("wrote", out)
