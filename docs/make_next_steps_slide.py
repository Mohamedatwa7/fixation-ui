"""Generate the F1X8 next-steps slide (16:9, noir theme) -> Downloads.
Usage: python docs/make_next_steps_slide.py
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

NOIR = RGBColor(0x0A, 0x0A, 0x0A)
PANEL = RGBColor(0x14, 0x14, 0x14)
ELEV = RGBColor(0x1C, 0x1C, 0x1C)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)
MUTED = RGBColor(0x9A, 0x9A, 0x9A)
ACCENT = RGBColor(0xE0, 0xE0, 0xE0)
GREEN = RGBColor(0x8A, 0xE1, 0xB8)
GOLD = RGBColor(0xFF, 0xC8, 0x61)

PROOF = [
    ("85/100", "image ranking accuracy — picks the better performer of two creatives (held-out Samsung posts)"),
    ("6×", "CTR spread measured between sibling KVs in one campaign (0.85% vs 0.14%, FF8 pre-order)"),
    ("~90 sec", "per image diagnostic; 3–6 min per video — vs weeks for a market test"),
    ("264", "Samsung creatives with real engagement labels already powering calibration"),
    ("0 lift", "from Meta's TRIBE brain model when we tested it on 171 posts — our behavioral stack already captures it"),
]

STEPS = [
    ("1. Turn on weekly retraining",
     "Accuracy grows with data: 60→142 training images moved ranking 0.838→0.851 AUC. Every week ~dozens of posts mature past the 14-day label window.",
     "Image ranking 85→87+ /100 pairwise, compounding weekly",
     "This week (built; one command to register)",
     "Scheduled task runs the existing retrain loop; a new model ships only if it beats the current one on holdout."),
    ("2. Train the reels ranker",
     "Video is the majority class (~719 labeled reels vs 265 images) yet still scored by the weaker method (73/100). Images jumped +16 pts when we switched judging→ranking.",
     "Video organic accuracy 73→85 /100 target; hook & watch-pull scoring precision",
     "2–3 weeks",
     "Same pairwise training recipe as images, on reel frames + motion, using engagement labels we already hold."),
    ("3. Add the paid-performance mode",
     "FF8 proved the organic lens can invert on paid offer ads: the 6×-CTR winner ranked last. Cost/engagement spread across siblings was 3.5× ($0.13 vs $0.45).",
     "New paid score calibrated to CTR & cost/engagement; CPM efficiency of creative selection",
     "As media exports land (~20–30 assets with CPM/CTR/ER)",
     "Ads-manager export (creative + results, matched by name) → calibrate a paid head next to the organic one."),
    ("4. Generalize across accounts",
     "Today's model part-learned 'Samsung Gulf aesthetics' (transfer accuracy 62/100 on a retailer feed). 12.5K non-brand posts are already exported and unused.",
     "Cross-account accuracy 62→75+ /100 — usable for any Samsung org or partner feed",
     "4–6 weeks, after steps 1–2",
     "Add non-brand extremes to ranker training; validate per-account before rollout."),
    ("5. Make pre-flight the workflow",
     "At hero reach (~20M impr.), the gap between a 0.2% and 0.9% ER creative is ~140K engagements — decided before a dirham is spent. Score → revise → rescore takes minutes.",
     "Campaign ER/CTR at launch; revision lift verified per asset; fewer wrong main-KV picks",
     "Immediately — process, not code",
     "Every KV/reel through f1x8.com before boosting: context in, In-Context score + brand-safe revision brief out."),
]

ADVANTAGE = ("The advantage: no competitor can copy the calibration — it is built on Samsung's own posts and results. "
             "Every campaign it scores makes it more accurate, and every point of ER it recovers is free media at launch scale.")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NOIR


def textbox(x, y, w, h, text, size, color, bold=False, font="Segoe UI",
            align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def cell_text(cell, text, size, color, bold=False, font="Segoe UI"):
    cell.margin_left = Emu(64008)
    cell.margin_right = Emu(64008)
    cell.margin_top = Emu(27432)
    cell.margin_bottom = Emu(27432)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color


# Header
textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.4),
        "F1X8 — NEXT STEPS", 12, MUTED, font="Consolas")
textbox(Inches(0.5), Inches(0.52), Inches(12.3), Inches(0.55),
        "From validated diagnostic to compounding creative advantage", 24, WHITE, bold=True)

# Proof band (where we are)
px = Inches(0.5)
pw = Inches(2.42)
for i, (num, label) in enumerate(PROOF):
    x = Emu(int(px) + i * (int(pw) + 91440 // 2))
    shape = slide.shapes.add_textbox(x, Inches(1.12), pw, Inches(0.86))
    tf = shape.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = num
    r1.font.size = Pt(17); r1.font.bold = True
    r1.font.name = "Consolas"; r1.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(8); r2.font.name = "Segoe UI"; r2.font.color.rgb = MUTED

# Steps table
top = Inches(2.12)
tbl_shape = slide.shapes.add_table(len(STEPS) + 1, 5, Inches(0.5), top,
                                   Inches(12.33), Inches(4.42))
tbl = tbl_shape.table
widths = [2.35, 3.55, 2.75, 1.60, 2.08]
for c, w in enumerate(widths):
    tbl.columns[c].width = Inches(w)

headers = ("WHAT", "WHY (the evidence)", "KPIs MOVED", "WHEN", "HOW")
for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ELEV
    cell_text(cell, h, 10, ACCENT, bold=True, font="Consolas")

for i, (what, why, kpi, when, how) in enumerate(STEPS, start=1):
    vals = [(what, WHITE, True, 9.5), (why, WHITE, False, 8.5),
            (kpi, GREEN, False, 8.5), (when, GOLD, False, 8.5),
            (how, MUTED, False, 8)]
    for c, (text, color, bold, size) in enumerate(vals):
        cell = tbl.cell(i, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if i % 2 else NOIR
        cell_text(cell, text, size, color, bold=bold)

# Advantage banner
textbox(Inches(0.5), Inches(6.72), Inches(12.33), Inches(0.65),
        ADVANTAGE, 12.5, ACCENT, bold=True, align=PP_ALIGN.CENTER)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Proof band sources: 85/100 = ranker holdout AUC 0.851 on unseen Samsung creatives; "
    "6x CTR = FF8 pre-order upper-funnel statics (q8h8offkv 0.85% vs b8kv 0.14%); 264 labeled "
    "creatives = calibration set (224 images + 40 reels); TRIBE test = eval/tribe/README.md, "
    "0 AUC lift on 171 posts, CC BY-NC research use.\n\n"
    "Step 3 detail: cost/engagement in FF8 lower statics ranged $0.13 (duo) to $0.45 (hero purple). "
    "Step 5 detail: 140K engagement gap = 20M impressions x (0.9% - 0.2%) ER.\n\n"
    "Positioning: F1X8 layers measured CV (OpenCV/TASED-Net gaze), LLM judgment (brand-playbook "
    "aware), and outcome calibration on Samsung's own engagement data. The moat is the data "
    "flywheel: each scored campaign adds labels, each label improves the ranker."
)

out = os.path.join(os.environ.get("USERPROFILE", ""), "Downloads", "F1X8_Next_Steps.pptx")
prs.save(out)
print("wrote", out)
