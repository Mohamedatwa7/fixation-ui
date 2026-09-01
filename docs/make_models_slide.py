"""Slide 2 for the stakeholder deck: the F1X8 AI models explained simply.
Matches the Samsung-internal light deck design (white bg, black section bars,
navy titles, blue/green/orange accents). -> Downloads/F1X8_Models_Slide.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1A, 0x1A, 0x2E)     # title / body text
GRAY = RGBColor(0x59, 0x59, 0x59)     # subtitle
LINE = RGBColor(0xBF, 0xBF, 0xBF)     # hairlines / box borders
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2F, 0x4B, 0xE0)     # keyword emphasis
GREEN = RGBColor(0x1E, 0x9E, 0x4A)    # positive metrics
ORANGE = RGBColor(0xE8, 0x86, 0x2B)   # process accents

# (header = model name, [(text, color, bold), ...] body parts, "covers" line)
# Ordered as the diagnostic story: see -> look -> listen -> measure -> compare -> judge.
CARDS = [
    ("Qwen2.5-VL",
     [("Visual understanding. ", NAVY, True),
      ("First, an AI ", NAVY, False),
      ("describes what is actually in the creative", BLUE, True),
      (" — products, faces, text, scenes — the way a human viewer would see it.", NAVY, False)],
     "Covers: Top risks · suggested fixes (the evidence behind them)"),
    ("TASED-Net",
     [("Attention mapping. ", NAVY, True),
      ("Then we predict where human eyes go in the first seconds — trained on ", NAVY, False),
      ("200K+ real eye-tracking recordings", BLUE, True),
      (" of people watching real content.", NAVY, False)],
     "Covers: Heatmaps · attention KPI scores"),
    ("Whisper + Librosa",
     [("Audio analysis. ", NAVY, True),
      ("For video, we listen too: ", NAVY, False),
      ("every spoken word transcribed", BLUE, True),
      (" and music energy tracked second by second, so sound and visuals are judged together.", NAVY, False)],
     "Covers: Video KPI scores · text translation"),
    ("OpenCV",
     [("Design measurement. ", NAVY, True),
      ("The layout is measured objectively — ", NAVY, False),
      ("hierarchy, contrast, white space, clutter", BLUE, True),
      (" — pure numbers, no opinion involved.", NAVY, False)],
     "Covers: KPI scores (the measured half)"),
    ("MAdVerse",
     [("Benchmark library. ", NAVY, True),
      ("Not a model — a reference set of ", NAVY, False),
      ("30K real ads measured the same way", BLUE, True),
      (", so every measurement becomes a percentile: better than X% of real ads.", NAVY, False)],
     "Covers: Benchmarks · the percentile behind every KPI"),
    ("Claude",
     [("Creative judge. ", NAVY, True),
      ("Finally, a senior-creative-director AI scores what can't be measured — emotion, brand strength, distinctiveness — ", NAVY, False),
      ("following the Samsung Brand Playbook", BLUE, True),
      (".", NAVY, False)],
     "Covers: Verdict · top risks · suggested fixes · image reiteration briefs"),
]

FOOTER_PARTS = [
    ("Every output on the previous slide traces back to one of these six engines — ", NAVY, False),
    ("measured and validated, not guessed", BLUE, True),
    (".", NAVY, False),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE


def runs_para(p, parts, size, font="Segoe UI"):
    for text, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color


# Title + subtitle + hairline (matches "F1X8 Model: ..." style)
tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.4), Inches(0.55))
p = tb.text_frame.paragraphs[0]
runs_para(p, [("F1X8 AI Models: What Measures What", NAVY, True)], 26)

ln = slide.shapes.add_shape(1, Inches(0.45), Inches(0.92), Inches(12.45), Emu(9525))
ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background()

tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.98), Inches(12.4), Inches(0.4))
p = tb.text_frame.paragraphs[0]
runs_para(p, [("The six engines behind the diagnosis — and which part of the output each one produces",
               GRAY, False)], 13)

# Card grid: 3 columns x 2 rows
COLS, ROWS = 3, 2
LEFT, TOP = Inches(0.45), Inches(1.50)
CARD_W, CARD_H = Inches(4.02), Inches(2.38)
GAP_X, GAP_Y = Inches(0.20), Inches(0.24)
BAR_H = Inches(0.38)

for i, (header, body_parts, powers) in enumerate(CARDS):
    row, col = i // COLS, i % COLS
    # Center a partial final row under the full rows above.
    in_row = min(len(CARDS) - row * COLS, COLS)
    row_indent = (COLS - in_row) * (int(CARD_W) + int(GAP_X)) // 2
    cx = Emu(int(LEFT) + row_indent + col * (int(CARD_W) + int(GAP_X)))
    cy = Emu(int(TOP) + row * (int(CARD_H) + int(GAP_Y)))

    # black header bar
    bar = slide.shapes.add_shape(1, cx, cy, CARD_W, BAR_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Emu(91440); tf.margin_top = Emu(18288); tf.margin_bottom = Emu(18288)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    runs_para(p, [(header, WHITE, True)], 13)

    # white body box with thin border
    body = slide.shapes.add_shape(1, cx, Emu(int(cy) + int(BAR_H)), CARD_W,
                                  Emu(int(CARD_H) - int(BAR_H)))
    body.fill.solid(); body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = LINE; body.line.width = Emu(9525)
    body.shadow.inherit = False
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(118872); tf.margin_right = Emu(118872); tf.margin_top = Emu(91440)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    runs_para(p, body_parts, 13.5)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(10)
    runs_para(p2, [("→ ", ORANGE, True), (powers, ORANGE, True)], 12.5)

# Footer line
tb = slide.shapes.add_textbox(Inches(0.45), Inches(6.62), Inches(12.45), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
runs_para(p, FOOTER_PARTS, 13)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Talk track: slide 1 showed the workflow (input -> AI models -> outputs). This slide opens the "
    "AI-models box. Six engines in story order (see, look, listen, measure, compare, judge), each with one job: TASED-Net = attention mapping trained on real "
    "human eye-tracking (the heatmaps); Qwen2.5-VL = visual understanding (what's in the frame); "
    "Whisper + librosa = audio transcription and sound-energy tracking; OpenCV = the measuring "
    "instrument for layout metrics, with MAdVerse (30K real ads) as the reference library that "
    "turns measurements into percentiles — instrument vs yardstick, they are different things; "
    "Claude = the creative judge scoring human dimensions under the Samsung Brand Playbook. "
    "Orange 'Covers' lines map each engine to the Diagnosis & Output column from slide 1. The "
    "fine-tuning/accuracy story (74%->85%) lives on the next slide."
)

out = os.path.join(os.environ.get("USERPROFILE", ""), "Downloads", "F1X8_Models_Slide.pptx")
prs.save(out)
print("wrote", out)
